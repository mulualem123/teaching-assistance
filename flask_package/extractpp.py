from pptx import Presentation
# Module for to get the title
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from . import googletransfun    #Orginal
from . import changealphabet    #Orginal


# class that take in power point and add the containts to the database. 
class extract:
    def __init__(self, powerpoint, db, filename):
        self.powerpoint = powerpoint
        self.db = db
        #filename is for the database entry. 
        self.filename = filename

    # A function to check if the files that has been selected are in the database
    ## You might have to change the id of the row.       
    def checkfile(self):
        exist = False
        files = self.db.get_data()
        for file in files:
            if file[5] == self.filename:
                exist = True
                print("extractpp checkfile file names from list" + str(file[5]))
                return exist
        print("extractpp checkfile filename recieved" + str(self.filename))    
        print("extractpp checkfile baloon" + str(exist))    
        return exist 
       
    #A function to read the containts of a power point and return the value
    def read(self):
         #read the power point
         self.powerpoint = Presentation(self.powerpoint)
         #get the slides
         slides = self.powerpoint.slides
         #loop through the slides
         ###loop through the slides and for each slides assign the Geez title, English title, Geez azmach, english alpha azmach, English translation
         ###
         slideTrack = 0
         for slide in slides:
             slideTrack +=1 #Troubleshoot tracks slide number.
             print ("Slide " + str(slideTrack)+"\n")
             trackTitle = 0 # tracks the number of Title. Makes sure is a title.
             trackEmpSpace = 0 # Make sure there is only few word in the title.
             track=0
             trackword = 0 #Troubleshoot-tracks the letters creating the tigle is for tracking a problem
             trackShape = 0 #Troubleshoot - tracks which shape is in the slide
             
             #get the shapes in the slide
             shapes = slide.shapes
             #initiate title
             geez_title_text="Mezmur"
             geez_text_content = ""
             en_apha_text_content = ""

             #loop through the shapes
             for shape in shapes:
                 trackShape += 1
                 print ("Shape " + str(trackShape)+"\n")
                 print (str(shape.shape_type))
                 #check if the shape ia a Title
                 if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    #print ("I am first" + str(track))
                    track+=1
                    placeholder = shape.placeholder_format
                    if placeholder.type == PP_PLACEHOLDER.TITLE:
                        # Now you have a title shape, and you can access its text
                        geez_title_text = shape.text
                        trackTitle += 1
                        #print(geez_title_text)
                    else:
                        geez_text_content += shape.text + "\n"
                 #check if the shape is a table
                 elif shape.has_table:
                     #print("I am second" + str(track))
                     track+=1
                     #get the table
                     table = shape.table
                     #loop through the rows
                     for row in table.rows:
                         #loop through the cells
                         for cell in row.cells:
                             #get the text in the cell
                             text = cell.text
                             #check if the text is not empty
                             if text:
                                 #add the text to the database
                                 self.db.mv_database(geez_title_text,text,self.filename)
                 else:
                    #print("I am last" + str(track))
                    track+=1
                    #check if the shape is a text box
                    if shape.has_text_frame:
                        #get the text frame
                        text_frame = shape.text_frame
                        #initiate the text paragraph
                        
                        #loop through the paragraphs in the text frame
                        for paragraph in text_frame.paragraphs:
                            #get the text in the paragraph
                            text = paragraph.text
                            #check if the text is not empty
                            if text:
                                #check language type
                                lgtype=googletransfun.check_language_type(text)
                                # checking if "text" is in english alphabet
                                if  lgtype== "am" or lgtype== "ti" :
                                    geez_text_content += paragraph.text + "\n"
                                    changal = changealphabet.geez_to_latin(text)
                                    print (str(lgtype) + " " + str(text))
                                    print (str(lgtype) + " " + str(changal))
                                    en_apha_text_content += changealphabet.geez_to_latin(text)
                                     # on the second space while going through the letters of text baloon to false and set geez_title_text=tit
                                    if trackTitle == 0:
                                        tit = " "
                                        for letter in text:
                                            if trackEmpSpace <=3:
                                                if letter == " ":
                                                    tit += letter
                                                    trackEmpSpace += 1
                                                else:
                                                    tit += letter
                                                    print ("Title in the loop" + str(tit))
                                                    trackword +=1 
                                                    print (trackword)          
                                        trackTitle +=1
                                    geez_title_text=tit                                                
                                else:
                                    en_apha_text_content = en_apha_text_content + text + "\n"
                                    print (str(lgtype) + " " + " else " + str(text))
                        #add the text to the database
                    else:
                        # #get the text frame
                        #text_frame = shape.text_frame
                        ##initiate the text paragraph
                        #
                        ##loop through the paragraphs in the text frame
                        #for paragraph in text_frame.paragraphs:
                        #    #get the text in the paragraph
                        #    text = paragraph.text
                        #    print("This is not in text frame. " + str(text))
                        print ("Thid is not shape.text_frame ")
             self.db.mv_database(geez_title_text, geez_text_content,en_apha_text_content,"NA",self.filename,"NA","NA","NA","NA")
## There might be a bug. more monitering required.
### Need each stence to be in a new line              
# set the title to the first alphabet of the paragraph and when track is zero --Solved
                                
                                 
            
 