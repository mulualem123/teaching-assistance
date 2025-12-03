import os
import base64
from datetime import datetime
#import requests
import requests # Uncommented: Required for Telegram send_message function
import fitz as fz
import sqlite3
from flask import Flask, render_template,send_from_directory, send_file, request as rq, flash, redirect, url_for, abort
from pptx import Presentation
from flask_mail import Mail, Message
from flask_sqlalchemy import SQLAlchemy #Account creation
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user, current_user #Account creation
#from flask_package.mailing import Gmail
from . import db                #Orginal
from .extractpp import extract  #Orginal
from . import googletransfun    #Orginal
from . import changealphabet    #Orginal
from .forms import LoginForm, PlaylistForm    #Orginal
from .models import db as account_db, User, Role, roles_users, Playlist, PlaylistSong, SavedFilter
from flask_security import Security, SQLAlchemyUserDatastore, UserMixin, RoleMixin, roles_required
from flask_migrate import Migrate  # Import Migrate here
from werkzeug.security import generate_password_hash
from werkzeug.utils import secure_filename
import click
from flask import jsonify
import json
import secrets

# Import the extended registration form
from .forms import ExtendedRegisterForm

app = Flask(__name__)
app.config["SECRET_KEY"]= b'\xa4\x99hM\x12s\xc3\x8d' # Moved to top: Best practice for app config

# Configure the upload folder for audio files
audio_folder = os.path.join(app.root_path, 'static', 'audio')
app.config['UPLOAD_FOLDER'] = audio_folder
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'pptx', 'ppt', 'mp3', 'wav', 'ogg', 'webm'}
# Ensure the upload folder exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# --- Path Definitions for Content Folders ---
# Define absolute paths for your content folders to avoid pathing issues.
pp_parent_folder = os.path.join(app.root_path, 'pp')
sundayClassPP = os.path.join(app.root_path, 'sundayClassPP')
os.makedirs(pp_parent_folder, exist_ok=True) # Ensure the directories exist
os.makedirs(sundayClassPP, exist_ok=True)

# --- Database Configuration ---
# Define an absolute path for the instance folder to ensure databases are always found.
instance_folder = os.path.join(app.root_path, 'instance')
os.makedirs(instance_folder, exist_ok=True)

# Configure the SQLite database for mezmur data (site.db)
app.config['DATABASE'] = os.path.join(instance_folder, 'site.db')

# Configure the SQLAlchemy database for user accounts (users.db)
app.config['SQLALCHEMY_DATABASE_URI'] = f"sqlite:///{os.path.join(app.instance_path, 'users.db')}"
app.config['SECURITY_PASSWORD_SALT'] = 'super-secret-salt'
# Tell Flask-Security to use our custom login and register templates
app.config['SECURITY_RECOVERABLE'] = True
app.config['SECURITY_CHANGEABLE'] = True
app.config['SECURITY_EMAIL_SENDER'] = 'no-reply@yourdomain.com'
app.config['SECURITY_REGISTERABLE'] = True # Enable the registration endpoint
app.config['SECURITY_SEND_REGISTER_EMAIL'] = False # Disable email sending on registration

app.config['SECURITY_LOGIN_USER_TEMPLATE'] = 'login.html'
app.config['SECURITY_REGISTER_USER_TEMPLATE'] = 'register.html'
app.config['SECURITY_REGISTER_FORM'] = ExtendedRegisterForm
# CSRF Configuration

db.init_app(app) #initiating sqlite3
print ("sql initiated")

account_db.init_app(app)    #Account creation
migrate = Migrate(app, account_db)
login_manager= LoginManager()  #Account creation
login_manager.init_app(app)  #Account creation
login_manager.login_view = 'login'  #Account creation

#Telegram
TOKEN = '7611669258:AAEchAugok05KQ_OqFzOc-59bY8FkSZQiwE'
TELEGRAM_API_URL = f'https://api.telegram.org/bot{TOKEN}/sendMessage'

# Setup Flask-Security
user_datastore = SQLAlchemyUserDatastore(account_db, User, Role)
security = Security(app, user_datastore)
        
with app.app_context():
    # Initialize both databases within the app context
    db.init_db()  # Creates site.db for mezmur data if it doesn't exist
    account_db.create_all() # Creates users.db for accounts if it doesn't exist
    print("Tables created successfully!")
    # Ensure saved_filters table has expected columns (migrate if necessary)
    try:
        users_db_path = os.path.join(app.instance_path, 'users.db')
        if os.path.exists(users_db_path):
            conn = sqlite3.connect(users_db_path)
            cur = conn.cursor()
            cur.execute("PRAGMA table_info(saved_filters)")
            cols = [row[1] for row in cur.fetchall()]
            # Add is_public column if missing
            if 'is_public' not in cols:
                try:
                    cur.execute("ALTER TABLE saved_filters ADD COLUMN is_public INTEGER DEFAULT 0")
                    print('Added is_public column to saved_filters')
                except Exception as e:
                    print('Failed to add is_public column:', e)
            # Add share_token column if missing
            if 'share_token' not in cols:
                try:
                    cur.execute("ALTER TABLE saved_filters ADD COLUMN share_token TEXT")
                    print('Added share_token column to saved_filters')
                except Exception as e:
                    print('Failed to add share_token column:', e)
            conn.commit()
            conn.close()
    except Exception as _err:
        print('Saved filters migration check failed:', _err)

# Import models after initializing account_db to avoid circular imports
@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

from .models import User    #Orginal
#gamail = Gmail(app)
#test = gamail.reminder  # test the connection with gmail server
app.config['MAIL_SERVER']='smtp.googlemail.com'
app.config['MAIL_PORT'] = 587
app.config['MAIL_USERNAME'] = 'mulualem.hailom@gmail.com'
app.config['MAIL_PASSWORD'] = ''
app.config['MAIL_USE_TLS'] = True
app.config['MAIL_USE_SSL'] = False
mail = Mail(app)
    
#Initialization 
geez_text = ""
#events to display on calendar card. Go to index page and call this array once you pass through home function
events = [
    {
        'title': 'event1',
        'date': '2024-10-15'
    },
    {
        'title': 'event2',
        'date': '2024-10-20'
    }
]

@app.cli.command("create-admin")
def create_admin_command():
    """Creates the default admin user from environment variables."""
    admin_email = os.getenv('ADMIN_EMAIL')
    admin_password = os.getenv('ADMIN_PASSWORD')

    if not admin_email or not admin_password:
        click.echo("Error: ADMIN_EMAIL and ADMIN_PASSWORD environment variables must be set.")
        return

    with app.app_context():
        if User.query.filter_by(email=admin_email).first():
            click.echo(f"Admin user with email '{admin_email}' already exists.")
            return

        try:
            admin_role = user_datastore.find_or_create_role(name='admin', description='Administrator')
            user_datastore.create_user(
                email=admin_email,
                username=admin_email.split('@')[0], # Added: Provide a username as it's a NOT NULL column
                password=admin_password, # Pass plaintext password to Flask-Security
                roles=[admin_role]
            )
            account_db.session.commit()
            click.echo(f"Admin user '{admin_email}' created successfully.")
        except Exception as e:
            account_db.session.rollback()
            click.echo(f"Error creating admin user: {e}")

@app.cli.command("reset-password")
@click.argument("email")
def reset_password_command(email):
    """Resets the password for a user with the given email."""
    new_password = os.getenv('NEW_PASSWORD')

    if not new_password:
        click.echo("Error: NEW_PASSWORD environment variable must be set.")
        return

    with app.app_context():
        user = User.query.filter_by(email=email).first()
        if not user:
            click.echo(f"User with email '{email}' not found.")
            return

        try:
            # Use Flask-Security's password hashing
            from flask_security.utils import hash_password
            user.password = hash_password(new_password)
            account_db.session.commit()
            click.echo(f"Password for '{email}' has been reset successfully.")
        except Exception as e:
            account_db.session.rollback()
            click.echo(f"Error resetting password: {e}")


    # Template helper to check for endpoint existence to avoid url_for build errors
    @app.context_processor
    def utility_processor():
        def endpoint_exists(name):
            try:
                return name in app.url_map._rules_by_endpoint
            except Exception:
                return False
        return dict(endpoint_exists=endpoint_exists)

@app.route('/register', methods=['GET', 'POST'])
def register():
    # This route is now handled by Flask-Security since we have set
    # SECURITY_REGISTER_USER_TEMPLATE. We can remove the custom logic.
    # Flask-Security will render the 'register.html' template automatically.
    return security.render_template('register.html')

        #user = User(username=form.username.data, email=form.email.data)
        #user.set_password(form.password.data)
        #account_db.session.add(user)
        
#Edit users
@app.route('/edit_user', methods=['GET', 'POST'])
@roles_required('admin')
def edit_user():
    if rq.method == 'POST':
        user_id = rq.form.get('user_id')
        
        user = User.query.get(user_id)
        if not user:
            return "User not found", 404
        
        #update user attribues
        new_username = rq.form.get('username')
        new_email = rq.form.get('email')
        new_active_status = rq.form.get('active')
        
        if new_username:    
            user.username = new_username
        else:
            return "User name can't be None", 400
        
        if new_email:
            user.email = new_email
        else:
            return "Email can't be None", 400
        
        if new_active_status:
            user.active = new_active_status.lower() == 'true'
        else:
            return "Active status can't be None", 400

        selected_roles = rq.form.getlist('roles[]')  # get the selected roles from the form
        user.roles = [] #Clear existing roles
        
        #Assign selected roles to the user
        if selected_roles:
            num = 1
            print ("This is going to print tags")
            for  role_name in selected_roles:
                user_datastore.add_role_to_user(user, role_name)
                print ("role " + str(num) + str(role_name))
                num = num + 1
        else:
            return "No roles selected", 400
        
        #Get selected action from the form
        selected_action = rq.form.get('user_action')
        
        print ("The selected action is " + str(selected_action))
        
        if selected_action.lower() == "activate":
            user.active = True
        elif selected_action.lower() == "deactivate":
            user.active = False
        elif selected_action.lower() == "delete":
            account_db.session.delete(user)
            
        #Add or remove roles
        #user_datastore.add_role_to_user(user, 'new_role')
        #user_datastore.remove_role_from_user(user, 'old_role')
        
        #Commit changes to the database
        account_db.session.commit()
        
        flash(str(user.username)+'\'s' + ' account has successfuly Edited!', 'success')
        return redirect(url_for('edit_user'))

    return render_template("user_manager.html",
        users= User.query.all(),
        roles= Role.query.all())
    
@app.route('/login', methods=['GET', 'POST'])
def login():
    # Flask-Security will automatically handle rendering the 'login.html' template
    # because we configured it in app.config. No need for custom logic here.
    return security.render_template('login.html')

@app.route('/favicon.ico')
def favicon():
    return send_from_directory(os.path.join(app.root_path, 'static'),
                               'logo.png', mimetype='image/png')

@app.route('/logout', methods=['GET', 'POST'])
@login_required
def logout():
    logout_user()
    flash('Logged out', 'warnning')
    return redirect(url_for('login'))

@app.route('/admin')
@roles_required('admin')
def admin():
    return render_template("user_manager.html",
            users= User.query.all(),
            roles= Role.query.all())

#Add role to the list 
@app.route('/add_role', methods=['GET', 'POST'])
def add_role():
    print ("add_role called in init.py")
    if rq.method == 'POST':
        role_name = rq.form.get('name')
        role_description = rq.form.get('description')
        user_datastore.create_role(name=str(role_name), description=str(role_name))
        account_db.session.commit()
        print (str(role_name) + str(role_description))
        for role in Role.query.all():
            print ("role: " + str(role) )
        return redirect(url_for('admin'))
    return render_template("user_manager.html",
        users= User.query.all(),
        roles= Role.query.all())
        
# A route to display a form for the user to enter a paragraph in Geez alphabet
@app.route("/")
def index():
    #account_db.create_all()
    #files = os.listdir(r'flask_package\pp')
    #return render_template("index.html", files=files, events=events)
    #files = os.listdir(r"C:\\Users\\selon\\Documents\\Bete Christian\\Mezmur")
    files = os.listdir(pp_parent_folder)
    imageslide1 = 'images/Emebatachn_Dingle_Mariam01.jpeg'
    imageslide2 = 'images/Pawlos01.jpeg'
    imageslide3 = 'images/images.jpeg'
    if geez_text != "": 
    #    return render_template("index.html", latin_text=changealphabet.geez_to_latin(my_map, geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
        return render_template("index.html", latin_text=changealphabet.geez_to_latin(geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data(),imageslide1=imageslide1, imageslide2=imageslide2,imageslide3=imageslide3)
    else:
        return render_template("index.html",files = os.listdir(pp_parent_folder), rows= db.get_data(),imageslide1=imageslide1, imageslide2=imageslide2,imageslide3=imageslide3)
    
# A route to display a form for the user to enter a paragraph in Geez alphabet
@app.route("/mezmur")
def mezmur():
    files = os.listdir(pp_parent_folder)
    form = PlaylistForm()
    shared_playlists = Playlist.query.filter_by(shared=True).order_by(Playlist.created_at.desc()).all()
    mez_tags = db.get_allMezTags()
    #print("The following will print all the mezmur tags")
    #for m_tags in mez_tags:
    #    print("m_tags " + str(m_tags))
    #    print("m_tags[0] " + str(m_tags[0]))
    #    print("m_tags[1] " + str(m_tags[1]))
    #    print("m_tags[2] " + str(m_tags[2]))
        
    #rows= db.get_data()  
    #audio_files={row[0]:audio_dic.get(row[0], None) for row in rows}
    
    return render_template("mezmur.html", 
                            latin_text=changealphabet.geez_to_latin(geez_text), 
                            lg_text = googletransfun.check_language_type(geez_text), 
                            geez_text_t = geez_text, 
                            translated_text = googletransfun.translate_tig_eng(geez_text),
                            files = os.listdir(pp_parent_folder), 
                            rows= db.get_data(),
                            mez_tags = mez_tags,
                            tags=db.get_taglist(),
                            form=form,
                            shared_playlists=shared_playlists
)


# Saved Filters API (list, create)
@app.route('/api/saved_filters', methods=['GET', 'POST'])
def api_saved_filters():
    # List saved filters for current user or create a new saved filter
    if rq.method == 'GET':
        if not current_user.is_authenticated:
            return jsonify([]), 401
        items = []
        for sf in account_db.session.query(SavedFilter).filter_by(user_id=current_user.id).order_by(SavedFilter.created_at.desc()).all():
            try:
                query_obj = json.loads(sf.query) if isinstance(sf.query, str) else sf.query
            except Exception:
                query_obj = {}
            items.append({
                'id': sf.id,
                'name': sf.name,
                'created_at': sf.created_at.isoformat() if sf.created_at else None,
                'query': query_obj,
                'is_public': bool(sf.is_public),
                'share_token': sf.share_token
            })
        return jsonify(items)

    # POST -> create
    if rq.method == 'POST':
        if not current_user.is_authenticated:
            return jsonify({'message': 'Authentication required'}), 401
        data = rq.get_json() or {}
        name = (data.get('name') or '').strip()
        if not name:
            return jsonify({'message': 'Name required'}), 400
        q = data.get('q', '')
        tags = data.get('tags', []) or []
        op = data.get('op', 'or')
        is_public = bool(data.get('is_public', False))

        # persist
        try:
            sf = SavedFilter(user_id=current_user.id, name=name, query=json.dumps({'q': q, 'tags': tags, 'op': op}), is_public=is_public)
            if is_public and not sf.share_token:
                sf.share_token = secrets.token_urlsafe(16)
            account_db.session.add(sf)
            account_db.session.commit()
            return jsonify({'id': sf.id, 'name': sf.name, 'created_at': sf.created_at.isoformat(), 'is_public': sf.is_public, 'share_token': sf.share_token}), 201
        except Exception as e:
            account_db.session.rollback()
            return jsonify({'message': 'Failed to save', 'error': str(e)}), 500


# Delete or update a saved filter
@app.route('/api/saved_filters/<int:sf_id>', methods=['DELETE', 'PUT'])
def api_saved_filter_item(sf_id):
    sf = account_db.session.get(SavedFilter, sf_id)
    if not sf:
        return jsonify({'message': 'Not found'}), 404

    # Ensure ownership
    if not current_user.is_authenticated or (sf.user_id != current_user.id and 'admin' not in [r.name for r in current_user.roles]):
        return jsonify({'message': 'Forbidden'}), 403

    if rq.method == 'DELETE':
        try:
            account_db.session.delete(sf)
            account_db.session.commit()
            return jsonify({'message': 'Deleted'}), 200
        except Exception as e:
            account_db.session.rollback()
            return jsonify({'message': 'Failed to delete', 'error': str(e)}), 500

    # PUT -> update is_public or name
    if rq.method == 'PUT':
        data = rq.get_json() or {}
        if 'is_public' in data:
            sf.is_public = bool(data.get('is_public'))
            if sf.is_public and not sf.share_token:
                sf.share_token = secrets.token_urlsafe(16)
        if 'name' in data:
            sf.name = (data.get('name') or sf.name).strip()
        try:
            account_db.session.add(sf)
            account_db.session.commit()
            return jsonify({'message': 'Updated'}), 200
        except Exception as e:
            account_db.session.rollback()
            return jsonify({'message': 'Failed to update', 'error': str(e)}), 500


@app.route("/sundayclass")
def sundayclass():
    files = os.listdir(sundayClassPP)
    
    #rows= db.get_data()  
    #audio_files={row[0]:audio_dic.get(row[0], None) for row in rows}
    
    if geez_text != "": 
        #return render_template("mezmur.html", latin_text=changealphabet.geez_to_latin(my_map, geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
        return render_template("sundayClass.html", latin_text=changealphabet.geez_to_latin(geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(sundayClassPP), rows= db.get_data())
    else:
        return render_template("sundayClass.html",files = os.listdir(sundayClassPP), rows= db.get_data())

# A route to display a form for the user to enter a paragraph in Geez alphabet
@app.route("/calendar")
def calendar():
    #files = os.listdir(r'flask_package\pp')
    #return render_template("index.html", files=files, events=events)
    #files = os.listdir(r"C:\\Users\\selon\\Documents\\Bete Christian\\Mezmur")
    files = os.listdir(pp_parent_folder)
    if geez_text != "": 
    #    return render_template("index.html", latin_text=changealphabet.geez_to_latin(my_map, geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
        return render_template("calendar.html", latin_text=changealphabet.geez_to_latin(geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
    else:
        return render_template("calendar.html",files = os.listdir(pp_parent_folder), rows= db.get_data())
    
@app.route("/files")
def files():
    #files = os.listdir(r'flask_package\pp')
    #return render_template("index.html", files=files, events=events)
    #files = os.listdir(r"C:\\Users\\selon\\Documents\\Bete Christian\\Mezmur")
    files = os.listdir(pp_parent_folder)
    if geez_text != "": 
    #    return render_template("index.html", latin_text=changealphabet.geez_to_latin(my_map, geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
        return render_template("files.html", latin_text=changealphabet.geez_to_latin(geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
    else:
        return render_template("files.html",files = os.listdir(pp_parent_folder), rows= db.get_data())
    
@app.route("/translate")
def translate():
    #files = os.listdir(r'flask_package\pp')
    #return render_template("index.html", files=files, events=events)
    #files = os.listdir(r"C:\\Users\\selon\\Documents\\Bete Christian\\Mezmur")
    files = os.listdir(pp_parent_folder)
    if geez_text != "": 
    #    return render_template("index.html", latin_text=changealphabet.geez_to_latin(my_map, geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
        return render_template("translate.html", 
                               latin_text=changealphabet.geez_to_latin(geez_text), 
                               lg_text = googletransfun.check_language_type(geez_text), 
                               geez_text_t = geez_text, 
                               translated_text = googletransfun.translate_tig_eng(geez_text),
                               files = os.listdir(pp_parent_folder), 
                               rows= db.get_data())
    else:
        return render_template("translate.html",files = os.listdir(pp_parent_folder), rows= db.get_data())
    
# Ensure 'requests' is imported if you're using it for Telegram
# import requests


#Return Audio with better error handling
@app.route("/audio/<id>")
def audio(id):
    try:
        # Get the audio file from the database using the passed id
        print("This is passed id from mezmur page: " + str(id))
        audio_list = db.get_audio(id)
        
        if audio_list and audio_list[0] and audio_list[0] != "NA":
            selected_audio = audio_list[0]
            print("001 audio/Selected_audio_exist " + str(selected_audio))
            
            # Construct the absolute path to the audio directory
            audio_dir = os.path.join(app.root_path, 'static', 'audio')
            
            # Check if file actually exists
            audio_path = os.path.join(audio_dir, str(selected_audio))
            if os.path.exists(audio_path):
                return send_from_directory(audio_dir, str(selected_audio))
            else:
                print(f"Audio file not found at path: {audio_path}")
                # Return a 404 error with proper JSON response
                return jsonify({
                    'error': 'Audio file not found',
                    'message': f'Audio file for ID {id} does not exist on server'
                }), 404
        else:
            print(f"002 No audio file found for ID: {id}")
            # Return a 404 error for missing audio
            return jsonify({
                'error': 'No audio available',
                    'message': f'No audio file associated with ID {id}' 
            }), 404
            
    except Exception as e:
        print(f"Error serving audio for ID {id}: {str(e)}")
        return jsonify({
            'error': 'Server error',
            'message': 'Failed to retrieve audio file'
        }), 500

#A function that returns text from given/selected PowerPoint file.
@app.route('/display/<filename>')
def display (filename):
    # Get the uploaded file from the request
    #f = pp_parent_folder + "\\" + filename
    f = pp_parent_folder + "/" + filename
    # Get the uploaded file from the request
    #f = str(r"C:\Users\selon\Documents\Bete Christian\Mezmur\1-29-2023.pptx")
    #print ("File name"  + str(f))
    # Get the file extension
    _, file_extension = os.path.splitext(f)
    #print ("File extension" + str(file_extension))
    
    # Initialize an empty string to store the extracted text
    text = ""
    # Check the file extension
    if file_extension == '.pptx':
        # Read the PowerPoint file using python-pptx
        prs = Presentation(f)
        ex = extract(f, db, filename)
        if ex.checkfile() != True:
            ex.read()
        else:
            print(str(filename) + " file exist!")
            
        # Loop over the slides and shapes
        for slide in prs.slides:
            #print("test2")
            for shape in slide.shapes:
                # Check if the shape has text
                if shape.has_text_frame and shape.text_frame.text:
                    # Append the text to the text string
                    text += shape.text + "\n"
                    #print(text)
    elif file_extension == '.pdf':
        # Read the PDF file using pymupdf
       # f = open(f, "w") # create a new file or overwrite an existing file
        doc = fz.open(filename=f, filetype="pdf")
        ex = extract(doc, db, filename)
        ex.read()
        # Loop over the pages
        for page in doc:
            # Get the text of the page
            text += page.get_text('text') + "\n"
    else:
        # Return an error message if the file extension is not supported
        return "Unsupported file format. Please upload a PowerPoint or PDF file."
    # Return the text as a Flask response
    #return render_template("index.html", text=text, lg_text = googletransfun.check_language_type(text),translated_text = googletransfun.translate_tig_eng(text), latin_text = changealphabet.geez_to_latin(my_map, text),files = os.listdir(r'C:\Users\selon\Documents\Bete Christian\Mezmur'))
    #return render_template("index.html", text=text, lg_text = googletransfun.check_language_type(text),translated_text = googletransfun.translate_tig_eng(text), latin_text = changealphabet.geez_to_latin(my_map, text),files = os.listdir(pp_parent_folder), rows= db.get_data())
    return render_template("index.html", text=text, lg_text = googletransfun.check_language_type(text),translated_text = googletransfun.translate_tig_eng(text), latin_text = changealphabet.geez_to_latin(text),files = os.listdir(pp_parent_folder), rows= db.get_data())

#This function downloads the files listed on the index page.
@app.route('/download/<filename>')
def download(filename):
    return send_from_directory(pp_parent_folder, filename, as_attachment=True)

# A route to process the user's input and display the converted paragraph in Latin alphabet
@app.route("/convert", methods=['GET', 'POST'])
def convert():
    geez_text = rq.form.get("geez_text")
    if rq.method == 'POST' and geez_text != "":
        title = "Mezmur"
        geez_text = rq.form.get("geez_text")
        #translated_text = googletransfun.translate_tig_eng(geez_text)
        #lg_text = googletransfun.check_language_type(geez_text)
        if 't_convert' in rq.form:
            #latin_text = changealphabet.geez_to_latin(my_map, geez_text)
            latin_text = changealphabet.geez_to_latin(geez_text)
            #db.mv_database(title,geez_text,latin_text,"NA","NA","NA","NA","NA")
            #mv_database(title,geez_text,latin_text,filename,audio,cat1,cat2,cat3):
            #rows = db.get_data()
            #print(f"latin_text '{latin_text}'")
            #return render_template("index.html", latin_text=latin_text, lg_text = lg_text, geez_text_t = geez_text, translated_text = translated_text,files = os.listdir(r'/python-docker/flask_package/pp'))
            #return render_template("index.html", latin_text=changealphabet.geez_to_latin(my_map, geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
            return render_template("translate.html", latin_text=changealphabet.geez_to_latin(geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
            #if geez_text != "": 
            #    return render_template("index.html", latin_text=changealphabet.geez_to_latin(my_map, geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
            #else:
            #    return render_template("index.html",files = os.listdir(pp_parent_folder), rows= db.get_data())            
        elif 'my_button' in rq.form:
            lg_text = googletransfun.translate_tig_eng()
            test = "Button clicked!"
            return render_template('translate.html', test_text=test, lg_text = lg_text)
    else:      
        msg = Message('Hello', sender = 'mulualem.hailom@gmail.com', recipients = ['hailomulalem@gmail.com'])
        msg.body = "Hello Flask message sent from Flask-Mail"
        mail.send(msg)
        flash("email have been sent successfuly")
        return render_template('index.html')
#Uppload File
#Check if the file extension is allowed
def allowed_file(filename):
    return '.' in filename and \
        filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']
        
@app.route("/uplaod", methods=['POST'])
def uplaod_file():
    if rq.method == 'POST':
        # check if the post request has the file part
        if 'file' not in rq.files:
            flash('No file part')
            return redirect(url_for('files'))

        file = rq.files['file']
        if file.filename == '':
            flash('No file selected for uploading')
            return redirect(url_for('files'))

        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            # Save the file to the 'pp' directory
            file.save(os.path.join(pp_parent_folder, filename))
            flash(f'File "{filename}" uploaded successfully!')
            return redirect(url_for('files'))
    return redirect(url_for('files'))

def upload (files):
    #check if the post request has the file part
    #if 'file' not in files:
    #    print(str(files))
    #    print("file is not in rq.files")
    #    return "No File exist"
    file = files['file']
    print(file)
    #Check if for empty file and has no name
    if file.filename == '':
        print(str(file.filename) + "the file has no name")
        return "File has no name"
    if file and allowed_file(file.filename):
        print("File does exit and is not null")
        filename = secure_filename(file.filename)
        file.save(os.path.join(app.config['UPLOAD_FOLDER'],filename))
        return filename
    else:
        return render_template("index.html",files = os.listdir(pp_parent_folder), rows= db.get_data())
  
@app.route('/delete/<id>')
@login_required
def delete (id):
    db.delete_data(id)
    #return render_template("index.html",files = os.listdir(pp_parent_folder), rows=db.get_data())
    form = PlaylistForm()

    return render_template("mezmur.html", 
                            latin_text=changealphabet.geez_to_latin(geez_text), 
                            lg_text = googletransfun.check_language_type(geez_text), 
                            geez_text_t = geez_text, 
                            translated_text = googletransfun.translate_tig_eng(geez_text),
                            files = os.listdir(pp_parent_folder), 
                            rows= db.get_data(),
                            mez_tags = db.get_allMezTags(),
                            tags=db.get_taglist(),
                            form = form
                           )
    
def generate_timed_lyrics(lyrics_text):
    if not lyrics_text:
        return ""
    lines = lyrics_text.strip().split('\n')
    timed_lines = [f"[00:00.000]{line.strip()}" for line in lines if line.strip()]
    return '\n'.join(timed_lines)

@app.route('/update/<id>',  methods=['GET', 'POST'])
@login_required
def update (id):
    #print(id)
    mezdata = db.get_selected_data(id) #Getting the selected row from db database
    print("id " + str(mezdata[0]))
    print("Title " + str(mezdata[1]))
    print("Titleen " + str(mezdata[2]))
    print("Azmach " + str(mezdata[3]))
    print("Azmachen" + str(mezdata[4]))
    print("number5 " + str(mezdata[5]))
    print("number6 " + str(mezdata[6]))
    print("number7 " + str(mezdata[7]))
    print("number8 " + str(mezdata[8]))
    print("number9 " + str(mezdata[9]))
    print("number10 " + str(mezdata[10]))
    print("number11 " + str(mezdata[11]))
    print("number12 " + str(mezdata[12]))
    print("number13 " + str(mezdata[13]))
    print("number14 " + str(mezdata[14]))
    

    #variable to be only passed or (updated and passed)
    latin_text = mezdata[4] 
    engTrans = mezdata[5]
    timed_geez = mezdata[6]
    timed_latin = mezdata[7]
    timed_english = mezdata[8]
    selected_mez_tags = db.get_selectedMezTags(id)
    
    for tag in selected_mez_tags:
        print ("This is selected mezmur's Tag " + str(tag))
       
    if mezdata[2] is None or mezdata[2] == " ":
        #db.set_titleen(changealphabet.geez_to_latin(my_map, mezdata[1]), id)
        db.set_titleen(changealphabet.geez_to_latin(mezdata[1]), id)
        #print ("Titleen ")
        #print (mezdata[2])
        
    if (mezdata[4]==None or mezdata[4]==""):
        print ("Azmach is empty or null")
        print ("The ID of the Mezmur is " + str(id))
        #db.set_azmachen(changealphabet.geez_to_latin(my_map, mezdata[3]), id)
        latin_text = changealphabet.geez_to_latin(mezdata[3])
        db.set_azmachen(latin_text, id)
        print("latin_text to be passed to update page is " + str(latin_text))

    if (mezdata[5]==None or mezdata[5]=="" or mezdata[5]=="NA"):
        print ("English translate is empty or null")
        engTrans = googletransfun.translate_tig_eng(mezdata[3])
        db.set_engTrans(engTrans,id)
        print ("English translate to be passed to update page is " + str(engTrans))

    if timed_geez is None or timed_geez.strip() == "":
        print("Timed Geez is empty or null, generating from Azmach")
        timed_geez = generate_timed_lyrics(mezdata[3])

    if timed_latin is None or timed_latin.strip() == "":
        print("Timed Latin is empty or null, generating from Azmachen")
        timed_latin = generate_timed_lyrics(latin_text)

    if timed_english is None or timed_english.strip() == "":
        print("Timed English is empty or null, generating from Hymn")
        timed_english = generate_timed_lyrics(engTrans)
    
    mezdata_list = list(mezdata)
    mezdata_list[6] = timed_geez
    mezdata_list[7] = timed_latin
    mezdata_list[8] = timed_english
        
    return render_template("update.html", 
                           mezmur=mezdata_list,
                           engTrat = googletransfun.check_language_type(mezdata[3]),
                           translated_text = googletransfun.translate_tig_eng(mezdata[3]),
                           latin_text = latin_text,
                           engTrans = engTrans,
                           tags=db.get_taglist(),
                           selected_mez_tags=selected_mez_tags)

@app.route('/pushupdate', methods=['POST'])
def pushupdate():
    # Configure allowed extensions and upload folder. Note that we've added 'webm' since recorded audio might be in this format.
    form = PlaylistForm()
    # Removed redundant app.config settings here; they are set globally at the top.

    if rq.method == 'POST':
        mezmur_id = rq.form.get("id")
        title = rq.form.get("title")
        titleen = rq.form.get("titleen")
        geez_text = rq.form.get("geez_text")
        alpha_text = rq.form.get("alpha_text")
        engTrans = rq.form.get("engTrans")
        timed_geez = rq.form.get("timed-geez")
        timed_latin = rq.form.get("timed-latin")
        timed_english = rq.form.get("timed-english")

        # Update the mezmur details in the database.
        db.set_title(title, mezmur_id)
        db.set_titleen(titleen, mezmur_id)
        db.set_azmach(geez_text, mezmur_id)
        db.set_azmachen(alpha_text, mezmur_id)
        db.set_engTrans(engTrans, mezmur_id)
        db.set_timed_geez(timed_geez, mezmur_id)
        db.set_timed_latin(timed_latin, mezmur_id)
        db.set_timed_english(timed_english, mezmur_id)

        # First, attempt to get an uploaded audio file.
        uploaded_filename = upload(rq.files)
        if "File has no name" not in uploaded_filename:
            db.set_audio_file(uploaded_filename, mezmur_id)
        else:
            # No file was uploaded. Check for recorded audio data.
            recorded_audio_data = rq.form.get("recorded_audio", "")
            if recorded_audio_data and recorded_audio_data.startswith("data:"):
                try:
                    # Split the data URL: e.g. "data:audio/webm; codecs=opus;base64,BASE64..."
                    header, encoded = recorded_audio_data.split(',', 1)
                    # Determine the file extension based on the header
                    if "audio/webm" in header:
                        extension = "webm"
                    elif "audio/ogg" in header:
                        extension = "ogg"
                    else:
                        extension = "mp3"  # Default to mp3 if type is unrecognized

                    filename = f"recorded_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{extension}"
                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)

                    # Write the decoded audio data to file.
                    with open(file_path, "wb") as f:
                        f.write(base64.b64decode(encoded))

                    db.set_audio_file(filename, mezmur_id)
                    print(f"Recorded audio saved as {filename}")
                except Exception as e:
                    print("Error saving recorded audio:", e)

        # Update selected tags (if any)
        selected_tags = rq.form.getlist('selected_tags')
        if selected_tags:
            print("Updating tags:")
            for num, tag in enumerate(selected_tags, start=1):
                print(f"Tag {num}: {tag}")
            db.update_mez_tags(selected_tags, mezmur_id)

    # Render the updated mezmur display page.
    return render_template("mezmur.html", 
                           latin_text=changealphabet.geez_to_latin(geez_text), 
                           lg_text=googletransfun.check_language_type(geez_text), 
                           geez_text_t=geez_text, 
                           translated_text=googletransfun.translate_tig_eng(geez_text),
                           files=os.listdir(pp_parent_folder), 
                           rows=db.get_data(),
                           mez_tags=db.get_allMezTags(),
                           tags=db.get_taglist(),
                           form=form)
    
@app.route('/add_mezmur', methods=['POST'])
def add_mezmur():
    if rq.method == 'POST':
        try:
            # ... your existing code to get form data ...
            title = rq.form.get("title")
            titleen = rq.form.get("titleen") # Added titleen
            geez_text = rq.form.get("geez_text")
            alpha_text = rq.form.get("alpha_text")
            engTrans = rq.form.get("engTrans")
            timed_geez = rq.form.get("timed-geez")
            timed_latin = rq.form.get("timed-latin")
            timed_english = rq.form.get("timed-english")
            selected_tags = rq.form.getlist('selected_tags') # Get selected tags
            
            # Handle file upload (as before)
            audioFile = rq.files.get("file")
            if audioFile and allowed_file(audioFile.filename):
                filename = secure_filename(audioFile.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                audioFile.save(filepath)
                audio_filepath = filepath
            else:
                audio_filepath = None

            # Call your database function
            result = db.add_mezmur(title, titleen, geez_text, alpha_text, engTrans, timed_geez, timed_latin, timed_english, audio_filepath, selected_tags)

            if isinstance(result, str) and result.startswith("Database error:"):
                return jsonify({'error': result}), 500
            else:
                return jsonify({'message': 'Mezmur added successfully'}), 201  # JSON success response

        except Exception as e:
            app.logger.exception(f"Error in add_mezmur: {e}")
            return jsonify({'error': 'An unexpected error occurred'}), 500

    return jsonify({'error': 'Invalid request method'}), 405  # Handle non-POST requests
        
@app.route('/selectedmez/<id>')
def selected(id):
    #print(id)
    text = db.get_selected_data(id)[3]
    #print (text)
    #return render_template("index.html", text=text, lg_text = googletransfun.check_language_type(text),translated_text = googletransfun.translate_tig_eng(text), latin_text = changealphabet.geez_to_latin(my_map, text),files = os.listdir(pp_parent_folder))  
    if geez_text != "": 
        #return render_template("index.html", latin_text=changealphabet.geez_to_latin(my_map, geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
        return render_template("index.html", latin_text=changealphabet.geez_to_latin(geez_text), lg_text = googletransfun.check_language_type(geez_text), geez_text_t = geez_text, translated_text = googletransfun.translate_tig_eng(geez_text),files = os.listdir(pp_parent_folder), rows= db.get_data())
    else:
        #return render_template("index.html", text=text, lg_text = googletransfun.check_language_type(text),translated_text = googletransfun.translate_tig_eng(text), latin_text = changealphabet.geez_to_latin(my_map, text),files = os.listdir(pp_parent_folder), rows= db.get_data())
        return render_template("index.html", text=text, lg_text = googletransfun.check_language_type(text),translated_text = googletransfun.translate_tig_eng(text), latin_text = changealphabet.geez_to_latin(text),files = os.listdir(pp_parent_folder), rows= db.get_data())
    
#Search function, takes any a word or a phrase and search in the data base
@app.route('/search', methods=['GET'])
def search():
    search_term = rq.args.get("search_term")
    if not search_term:
        #return "No search term provided", 400
        return  redirect(url_for('mezmur'))
    
    #Search the database for the rows where the title or the ext column contains the term
    results = db.search(search_term)
    
    #Render the sarch results template, passing in the search term and results
    return render_template("mezmur.html",files = os.listdir(pp_parent_folder), search_term=search_term, rows= results)


#Add tag to list through form
@app.route('/add_tag_form', methods=['GET', 'POST'])
def add_tag_form():
    # This route is primarily for displaying the tag management page.
    # POST requests for adding tags are now handled by the /add_tag API endpoint.
    return render_template("add_tag_form.html", tags=db.get_taglist())

# API endpoint for adding a tag (used by AJAX from add_tag_form.html)
@app.route('/add_tag', methods=['POST'])
def add_tag():
    if rq.method == 'POST':
        tag_name = rq.form.get('tag_name').strip()
        if not tag_name:
            return jsonify({'error': 'Tag name cannot be empty'}), 400

        # Check if tag already exists (case-insensitive)
        if db.tag_exists(tag_name):
            return jsonify({'error': 'Tag already exists'}), 409

        try:
            new_tag_id = db.add_tag(tag_name) # db.add_tag now returns the new ID
            return jsonify({'message': 'Tag added successfully', 'tag': tag_name, 'tag_id': new_tag_id}), 201
        except Exception as e:
            app.logger.error(f"Error adding tag: {e}")
            return jsonify({'error': 'Failed to add tag'}), 500

    return jsonify({'error': 'Invalid request method'}), 405
 
#Delete_tag  
@app.route('/delete_tag/<int:tag_id>', methods=['POST'])
@login_required
@roles_required('admin', 'mezmur_editor') # Only admins or mezmur editors can delete tags
def delete_tag_api(tag_id):
    try:
        rows_affected = db.delete_tag(tag_id)
        if rows_affected > 0:
            return jsonify({'message': 'Tag deleted successfully'}), 200
        else:
            return jsonify({'error': 'Tag not found'}), 404
    except Exception as e:
        app.logger.error(f"Error deleting tag {tag_id}: {e}")
        return jsonify({'error': 'Failed to delete tag'}), 500

# API endpoint for updating a tag (used by AJAX from add_tag_form.html)
@app.route('/update_tag/<int:tag_id>', methods=['POST'])
@login_required
@roles_required('admin', 'mezmur_editor') # Only admins or mezmur editors can update tags
def update_tag_api(tag_id):
    try:
        data = rq.get_json()
        new_name = data.get('new_name', '').strip()
        if not new_name:
            return jsonify({'error': 'New tag name cannot be empty'}), 400
        db.update_tag(tag_id, new_name)
        return jsonify({'message': 'Tag updated successfully', 'tag': new_name}), 200
    except Exception as e:
        app.logger.error(f"Error updating tag {tag_id}: {e}")
        return jsonify({'error': 'Failed to update tag'}), 500

#Attach a Tag to a Mezmur
@app.route('/add_tag_to_mezmur/<int:mez_id>', methods=['GET', 'POST'])
def add_tag_to_mez(mez_id):
    tag_name = rq.form.get('name')
    db.add_tags_to_mez(mez_id,tag_name)
    return redirect(url_for('mezmur'))

#getting built #Telegram
@app.route('/send_message', methods=['POST'])
def send_message():
    data = rq.json
    chat_id = data.get('chat_id')
    message = data.get('message')
    
    payload = {
        'chat_id': chat_id,
        'text': message
    }
    
    response = rq.post(TELEGRAM_API_URL, json=payload)
    return response.json()
################################PlayList###################
# Create a playlist
@app.route('/playlist/create', methods=['GET', 'POST'])
@login_required
def create_playlist():
    form = PlaylistForm()
    if form.validate_on_submit():
        new_playlist = Playlist(user_id=current_user.id, name=form.name.data, description=form.description.data)
        account_db.session.add(new_playlist)
        account_db.session.commit()
        flash('playlist created!', 'success')
        return redirect(url_for('mezmur'))
    return render_template('create_playlist.html', form=form)

# View all playlists
@app.route('/playlists')
@login_required
def view_playlists():
    playlists = Playlist.query.filter_by(user_id=current_user.id).all()
    return render_template('playlists.html', playlists=playlists)

#Gets list of playlists with caching and optimization
@app.route('/api/playlists')
@login_required
def get_playlists():
    try:
        # Add caching headers for better performance
        playlists = Playlist.query.filter_by(user_id=current_user.id).all()
        
        response_data = []
        for p in playlists:
            # More efficient song count query
            song_count = account_db.session.query(PlaylistSong).filter_by(playlist_id=p.id).count()
            response_data.append({
                'id': p.id,
                'name': p.name,
                'description': p.description or '',
                'song_count': song_count,
                'shared': p.shared if hasattr(p, 'shared') else False,
                'created_at': p.created_at.isoformat() if hasattr(p, 'created_at') else None
            })
        
        response = jsonify(response_data)
        # Add caching headers
        response.headers['Cache-Control'] = 'private, max-age=300'  # 5 minutes
        return response
    except Exception as e:
        app.logger.error(f"Error fetching playlists: {str(e)}")
        return jsonify({
            'status': 'error',
            'message': 'Failed to fetch playlists',
            'details': str(e) if app.debug else 'Internal server error'
        }), 500
        
#Load songs from selected playlist with optimization
@app.route('/api/playlists/<int:playlist_id>')
@login_required
def get_playlist(playlist_id):
    try:
        # Check if user is admin
        is_admin = 'admin' in [role.name for role in current_user.roles]
        
        # Get playlist from database with better error handling
        if is_admin:
            # Admins can access any playlist
            playlist = Playlist.query.get(playlist_id)
        else:
            # Regular users can access their own playlists or shared playlists
            playlist = Playlist.query.get(playlist_id)
            
            # Check access permissions for non-admin users
            if playlist and playlist.user_id != current_user.id and not playlist.shared:
                # User is not owner and playlist is not shared - deny access
                return jsonify({
                    'status': 'error',
                    'message': 'Playlist not found or access denied'
                }), 404
        
        if not playlist:
            return jsonify({
                'status': 'error',
                'message': 'Playlist not found or access denied'
            }), 404
        
        # If it's a shared playlist from another user, use the shared playlist data format
        if playlist.user_id != current_user.id and playlist.shared:
            return get_shared_playlist_data(playlist)

        # Optimize song loading with batch queries              
        songs = []
        song_relations = playlist.songs.all()
        
        # Batch load song data to reduce database calls
        song_ids = [rel.song_id for rel in song_relations]
        
        for song_id in song_ids:
            song_data = db.get_selected_data(song_id)
            if song_data:
                songs.append({
                    'id': song_data[0],
                    'title': song_data[1],
                    'timed_geez': song_data[6] if len(song_data) > 6 else '',
                    'timed_latin': song_data[7] if len(song_data) > 7 else '',
                    'timed_english': song_data[8] if len(song_data) > 8 else '',
                    'azmach': song_data[3] if len(song_data) > 3 else '',
                    'azmachen': song_data[4] if len(song_data) > 4 else '',
                    'engTrans': song_data[5] if len(song_data) > 5 else '',
                    'audio_url': url_for('audio', id=song_data[0], _external=True)
                })

        response_data = {
            'id': playlist.id,
            'name': playlist.name,
            'description': playlist.description or '',
            'songs': songs,
            'total_songs': len(songs),
            'user_id': playlist.user_id,
            'shared': playlist.shared if hasattr(playlist, 'shared') else False
        }
        
        response = jsonify(response_data)
        # Add appropriate caching headers
        response.headers['Cache-Control'] = 'private, max-age=60'  # 1 minute for playlist content
        return response

    except Exception as e:
        app.logger.error(f"Error fetching playlist {playlist_id}: {str(e)}")
        return jsonify({
            'status': 'error',
            'message': 'Failed to fetch playlist',
            'details': str(e) if app.debug else 'Internal server error'
        }), 500

# Public playlist API (no login required for shared playlists)
@app.route('/api/playlists/public/<int:playlist_id>')
def get_public_playlist(playlist_id):
    try:
        # Only allow access to shared playlists
        playlist = Playlist.query.filter_by(
            id=playlist_id,
            shared=True
        ).first()
        
        if not playlist:
            return jsonify({
                'status': 'error',
                'message': 'Public playlist not found'
            }), 404

        return get_shared_playlist_data(playlist)

    except Exception as e:
        app.logger.error(f"Error fetching public playlist {playlist_id}: {str(e)}")
        return jsonify({
            'status': 'error',
            'message': 'Failed to fetch playlist',
            'details': str(e) if app.debug else 'Internal server error'
        }), 500

def get_shared_playlist_data(playlist):
    """Helper function to format playlist data for public access"""
    songs = []
    song_relations = playlist.songs.all()
    
    # Batch load song data to reduce database calls
    song_ids = [rel.song_id for rel in song_relations]
    
    for song_id in song_ids:
        song_data = db.get_selected_data(song_id)
        if song_data:
            songs.append({
                'id': song_data[0],
                'title': song_data[1],
                'timed_geez': song_data[6] if len(song_data) > 6 else '',
                'timed_latin': song_data[7] if len(song_data) > 7 else '',
                'timed_english': song_data[8] if len(song_data) > 8 else '',
                'azmach': song_data[3] if len(song_data) > 3 else '',
                'azmachen': song_data[4] if len(song_data) > 4 else '',
                'engTrans': song_data[5] if len(song_data) > 5 else '',
                'audio_url': url_for('audio', id=song_data[0], _external=True)
            })

    response_data = {
        'id': playlist.id,
        'name': playlist.name,
        'description': playlist.description or '',
        'songs': songs,
        'total_songs': len(songs),
        'shared': True,
        'owner': playlist.user.username if playlist.user else 'Unknown'
    }
    
    response = jsonify(response_data)
    # Add appropriate caching headers for public content
    response.headers['Cache-Control'] = 'public, max-age=300'  # 5 minutes for public content
    return response

# Get individual song details API
@app.route('/api/song/<int:song_id>')
def get_song_details(song_id):
    try:
        song_data = db.get_selected_data(song_id)
        
        if not song_data:
            return jsonify({
                'status': 'error',
                'message': 'Song not found'
            }), 404
        
        response_data = {
            'status': 'success',
            'song': {
                'id': song_data[0],
                'title': song_data[1],
                'title_en': song_data[2] if len(song_data) > 2 else '',
                'geez_text': song_data[3] if len(song_data) > 3 else '',
                'latin_text': song_data[4] if len(song_data) > 4 else '',
                'english_text': song_data[5] if len(song_data) > 5 else '',
                'timed_geez': song_data[6] if len(song_data) > 6 else '',
                'timed_latin': song_data[7] if len(song_data) > 7 else '',
                'timed_english': song_data[8] if len(song_data) > 8 else '',
                'audio_url': url_for('audio', id=song_data[0], _external=True)
            }
        }
        
        response = jsonify(response_data)
        response.headers['Cache-Control'] = 'public, max-age=600'  # 10 minutes for song data
        return response
        
    except Exception as e:
        app.logger.error(f"Error fetching song {song_id}: {str(e)}")
        return jsonify({
            'status': 'error',
            'message': 'Failed to fetch song details',
            'details': str(e) if app.debug else 'Internal server error'
        }), 500

# Debug route to check playlist status (remove in production)
@app.route('/api/debug/playlists')
def debug_playlists():
    """Debug endpoint to check all playlists - remove in production"""
    if not app.debug:
        return jsonify({'error': 'Debug mode only'}), 403
    
    try:
        all_playlists = Playlist.query.all()
        debug_info = []
        
        for playlist in all_playlists:
            debug_info.append({
                'id': playlist.id,
                'name': playlist.name,
                'user_id': playlist.user_id,
                'shared': getattr(playlist, 'shared', False),
                'song_count': playlist.songs.count(),
                'owner': playlist.user.username if playlist.user else 'No owner'
            })
        
        return jsonify({
            'total_playlists': len(all_playlists),
            'playlists': debug_info,
            'current_user': {
                'id': current_user.id if current_user.is_authenticated else None,
                'username': current_user.username if current_user.is_authenticated else 'Not logged in'
            }
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# API endpoint to get available tags
@app.route('/api/tags')
def get_tags():
    """Get all available tags with counts"""
    try:
        tags = db.get_taglist()  # Using your existing function
        
        # Format tags with counts if your db function supports it
        if tags:
            formatted_tags = []
            for tag in tags:
                if isinstance(tag, tuple):
                    # If tag includes count: (name, count)
                    formatted_tags.append({
                        'name': tag[0],
                        'count': tag[1] if len(tag) > 1 else 0
                    })
                else:
                    # If tag is just a string
                    formatted_tags.append({
                        'name': str(tag),
                        'count': 0
                    })
            return jsonify(formatted_tags)
        else:
            return jsonify([])
            
    except Exception as e:
        app.logger.error(f"Error fetching tags: {str(e)}")
        return jsonify({'error': 'Failed to fetch tags'}), 500


    @app.route('/api/saved_filters', methods=['GET'])
    @login_required
    def get_saved_filters():
        try:
            saved = SavedFilter.query.filter_by(user_id=current_user.id).order_by(SavedFilter.created_at.desc()).all()
            result = []
            for s in saved:
                try:
                    payload = json.loads(s.query)
                except Exception:
                    payload = {'q': '', 'tags': [], 'op': 'or'}

                result.append({
                    'id': s.id,
                    'name': s.name,
                    'query': payload,
                    'created_at': s.created_at.isoformat() if s.created_at else None
                })

            return jsonify(result)
        except Exception as e:
            app.logger.exception(f"Error in get_saved_filters: {e}")
            return jsonify({'error': 'Failed to load saved filters'}), 500


    @app.route('/profile')
    @login_required
    def profile():
        """User profile page - lets user manage saved filters"""
        return render_template('profile.html')


    @app.route('/api/saved_filters', methods=['POST'])
    @login_required
    def create_saved_filter():
        try:
            data = rq.get_json() or {}
            name = data.get('name', '').strip()
            if not name:
                return jsonify({'error': 'Filter name is required'}), 400

            q = data.get('q', '')
            tags = data.get('tags', [])
            op = data.get('op', 'or')

            payload = json.dumps({'q': q, 'tags': tags, 'op': op})

            # allow public flag optionally
            is_public = bool(data.get('is_public', False))
            share_token = None
            if is_public:
                # create a small UUID-like token
                import secrets
                share_token = secrets.token_urlsafe(16)

            sf = SavedFilter(user_id=current_user.id, name=name, query=payload, is_public=is_public, share_token=share_token)
            account_db.session.add(sf)
            account_db.session.commit()

            return jsonify({'id': sf.id, 'name': sf.name, 'query': json.loads(sf.query)}), 201
        except Exception as e:
            app.logger.exception(f"Error creating saved filter: {e}")
            return jsonify({'error': 'Failed to save filter'}), 500


    @app.route('/api/saved_filters/<int:filter_id>', methods=['DELETE'])
    @login_required
    def delete_saved_filter(filter_id):
        try:
            sf = SavedFilter.query.get(filter_id)
            if not sf:
                return jsonify({'error': 'Not found'}), 404
            if sf.user_id != current_user.id:
                return jsonify({'error': 'Forbidden'}), 403

            account_db.session.delete(sf)
            account_db.session.commit()
            return jsonify({'status': 'deleted'})
        except Exception as e:
            app.logger.exception(f"Error deleting saved filter: {e}")
            return jsonify({'error': 'Failed to delete filter'}), 500


    @app.route('/api/saved_filters/<int:filter_id>', methods=['PUT'])
    @login_required
    def update_saved_filter(filter_id):
        try:
            sf = SavedFilter.query.get(filter_id)
            if not sf:
                return jsonify({'error': 'Not found'}), 404
            if sf.user_id != current_user.id:
                return jsonify({'error': 'Forbidden'}), 403

            data = rq.get_json() or {}
            # allow toggling public
            if 'is_public' in data:
                is_public = bool(data.get('is_public'))
                sf.is_public = is_public
                if is_public and not sf.share_token:
                    import secrets
                    sf.share_token = secrets.token_urlsafe(16)
                if not is_public:
                    sf.share_token = None

            account_db.session.commit()
            return jsonify({'id': sf.id, 'is_public': sf.is_public, 'share_token': sf.share_token})
        except Exception as e:
            app.logger.exception(f"Error updating saved filter: {e}")
            return jsonify({'error': 'Failed to update saved filter'}), 500


    @app.route('/api/public_filters', methods=['GET'])
    def list_public_filters():
        try:
            public = SavedFilter.query.filter_by(is_public=True).order_by(SavedFilter.created_at.desc()).all()
            return jsonify([{
                'id': s.id,
                'name': s.name,
                'query': json.loads(s.query),
                'share_token': s.share_token,
                'created_at': s.created_at.isoformat() if s.created_at else None,
                'user_id': s.user_id
            } for s in public])
        except Exception as e:
            app.logger.exception(f"Error listing public filters: {e}")
            return jsonify({'error':'failed'}), 500


    @app.route('/api/saved_filters/shared/<token>', methods=['GET'])
    def get_shared_filter(token):
        try:
            sf = SavedFilter.query.filter_by(share_token=token, is_public=True).first()
            if not sf:
                return jsonify({'error': 'Not found'}), 404
            return jsonify({'id': sf.id, 'name': sf.name, 'query': json.loads(sf.query), 'user_id': sf.user_id})
        except Exception as e:
            app.logger.exception(f"Error retrieving shared filter: {e}")
            return jsonify({'error':'failed'}), 500


@app.route('/api/mezmurs')
def api_mezmurs():
    """API: GET /api/mezmurs?q=&tags=a,b&op=and&page=1&perPage=24
    Returns JSON: { total, page, perPage, items: [{m_id,title,azmach,engTrans,dir,created}], facets: {tags:[{name,count}]}}
    """
    try:
        q = rq.args.get('q', '').strip()
        tags_param = rq.args.get('tags', '').strip()
        tags = [t.strip() for t in tags_param.split(',') if t.strip()]
        op = rq.args.get('op', 'or').lower()
        page = max(int(rq.args.get('page', 1)), 1)
        per_page = max(int(rq.args.get('perPage', 24)), 1)

        conn = db.get_db()
        cursor = conn.cursor()

        where_clauses = []
        params = []

        if q:
            # Search title, azmach, azmachen, engTrans
            like = f"%{q}%"
            where_clauses.append("(title LIKE ? OR azmach LIKE ? OR azmachen LIKE ? OR engTrans LIKE ?)")
            params.extend([like, like, like, like])

        if tags:
            if op == 'and':
                # require each tag exists for a mezmur
                for t in tags:
                    where_clauses.append("EXISTS (SELECT 1 FROM mezTags mt WHERE mt.m_id = mezmur.m_id AND mt.tag = ?)")
                    params.append(t)
            else:
                # OR semantics - mezmur appears if any of the tags match
                placeholders = ','.join(['?'] * len(tags))
                where_clauses.append(f"mezmur.m_id IN (SELECT m_id FROM mezTags mt WHERE mt.tag IN ({placeholders}))")
                params.extend(tags)

        where_sql = ('WHERE ' + ' AND '.join(where_clauses)) if where_clauses else ''

        # total count
        count_sql = f"SELECT COUNT(*) as total FROM mezmur {where_sql}"
        cursor.execute(count_sql, params)
        total = cursor.fetchone()[0] or 0

        # items - basic columns
        offset = (page - 1) * per_page
        items_sql = f"SELECT m_id, title, titleen, azmach, azmachen, engTrans, dir, audio_file, created FROM mezmur {where_sql} ORDER BY m_id DESC LIMIT ? OFFSET ?"
        items_params = params + [per_page, offset]
        cursor.execute(items_sql, items_params)
        rows = cursor.fetchall()

        items = []
        for r in rows:
            items.append({
                'm_id': r['m_id'],
                'title': r['title'],
                'titleen': r['titleen'],
                'azmach': r['azmach'],
                'azmachen': r['azmachen'],
                'engTrans': r['engTrans'],
                'dir': r['dir'],
                'audio_file': r['audio_file'],
                'created': r['created']
            })

        # Facet counts -> compute counts per tag that reflect the current search q
        # For each tag t, compute number of mezmur that would match when t is included along with the current q
        tag_list = [r['tag'] for r in conn.execute('SELECT tag FROM tagList').fetchall()]

        facets = []
        # Build q-only where clause to use as base
        base_where = ""
        base_params = []
        if q:
            like = f"%{q}%"
            base_where = "WHERE (title LIKE ? OR azmach LIKE ? OR azmachen LIKE ? OR engTrans LIKE ? )"
            base_params = [like, like, like, like]

        # If there are no selected tags, we can compute simple counts by joining base query with mezTags
        for t in tag_list:
            try:
                # effective tags set = current tags + candidate t (ensure unique)
                effective_tags = list(dict.fromkeys(tags + [t]))

                if effective_tags:
                    # For OR: any of effective_tags
                    if op == 'or':
                        placeholders = ','.join(['?'] * len(effective_tags))
                        facet_sql = f"SELECT COUNT(DISTINCT mezmur.m_id) as cnt FROM mezmur JOIN mezTags mt ON mt.m_id = mezmur.m_id {base_where} AND mezmur.m_id IN (SELECT m_id FROM mezTags WHERE tag IN ({placeholders}))"
                        params_for_facet = base_params + effective_tags
                    else:
                        # AND: mezmur must have all effective_tags
                        # Build EXISTS clauses for each tag
                        exists_clauses = ' AND '.join(["EXISTS (SELECT 1 FROM mezTags mt2 WHERE mt2.m_id = mezmur.m_id AND mt2.tag = ?)" for _ in effective_tags])
                        facet_sql = f"SELECT COUNT(*) as cnt FROM mezmur {base_where} AND {exists_clauses}"
                        params_for_facet = base_params + effective_tags
                else:
                    # No tag filter -> count all matching base (q)
                    facet_sql = f"SELECT COUNT(*) as cnt FROM mezmur {base_where}"
                    params_for_facet = base_params

                cursor.execute(facet_sql, params_for_facet)
                cnt = cursor.fetchone()[0] or 0
                facets.append({'name': t, 'count': cnt})
            except Exception:
                # On any failure just return 0 for safety
                facets.append({'name': t, 'count': 0})

        # sort facets by count desc to return useful ordering
        facets.sort(key=lambda x: x['count'], reverse=True)

        return jsonify({
            'total': total,
            'page': page,
            'perPage': per_page,
            'items': items,
            'facets': {
                'tags': facets
            }
        })

    except Exception as e:
        app.logger.exception(f"Error in /api/mezmurs: {e}")
        return jsonify({'error': 'Failed to fetch mezmurs'}), 500
        
# View a single playlist and its songs
@app.route('/playlist/<int:playlist_id>')
@login_required
def view_playlist(playlist_id):
    playlist = Playlist.query.get_or_404(playlist_id)
    if playlist.user_id != current_user.id:
        flash('Access denied!', 'danger')
        return redirect(url_for('mezmur'))
    
    # Fetch song details from SQLite (`db`)
    songs = []
    for song_rel in playlist.songs.all():
        song_data = db.get_selected_data(song_rel.song_id)  # Use your SQLite helper
        if song_data:
            songs.append(song_data)
    
    return render_template('playlist.html', playlist=playlist, songs=songs)

@app.route('/playlist/<int:playlist_id>/delete', methods=['POST'])
@login_required
def delete_playlist_old(playlist_id):
    playlist = Playlist.query.get_or_404(playlist_id)
    if playlist.user_id != current_user.id:
        abort(403)
    # Delete associated songs
    PlaylistSong.query.filter_by(playlist_id=playlist_id).delete()
    account_db.session.delete(playlist)
    account_db.session.commit()
    flash('Playlist deleted successfully!', 'success')
    return redirect(url_for('view_playlists'))

@app.route('/api/playlists/<int:playlist_id>', methods=['DELETE'])
@login_required
def delete_playlist(playlist_id):
    try:
        # Get playlist
        playlist = Playlist.query.get_or_404(playlist_id)

        # Verify ownership or admin role
        is_admin = any(role.name == 'admin' for role in current_user.roles)
        if playlist.user_id != current_user.id and not is_admin:
            return jsonify({'status': 'error', 'message': 'Permission denied'}), 403

        # Delete all associated songs first
        PlaylistSong.query.filter_by(playlist_id=playlist_id).delete()
        
        # Delete the playlist
        account_db.session.delete(playlist)
        account_db.session.commit()

        # Get updated list of playlists
        updated_playlists = Playlist.query.filter_by(
            user_id=current_user.id
        ).all()

        return jsonify({
            'status': 'success',
            'message': 'Playlist deleted successfully',
            'playlists': [{
                'id': p.id,
                'name': p.name,
                'description': p.description,
                'song_count': p.songs.count()
            } for p in updated_playlists]
        })

    except Exception as e:
        account_db.session.rollback()
        return jsonify({
            'status': 'error',
            'message': str(e)
        }), 500

@app.route('/api/playlist/<int:playlist_id>/share', methods=['POST'])
@login_required
def share_playlist(playlist_id):
    playlist = Playlist.query.get_or_404(playlist_id)

    # Check ownership
    if playlist.user_id != current_user.id:
        return jsonify({'status': 'error', 'message': 'Permission denied'}), 403

    # Toggle the shared status
    playlist.shared = not playlist.shared
    account_db.session.commit()

    return jsonify({
        'status': 'success',
        'message': f'Playlist is now {"public" if playlist.shared else "private"}.',
        'shared': playlist.shared
    })
           
# Add a song to a playlist
@app.route('/add_to_playlist', methods=['POST'])
@login_required
def add_to_playlist_form():
    song_id = rq.form.get('song_id')
    playlist_id = rq.form.get('playlist_id')
    
    # Validate song exists in SQLite (`db`)
    if not db.get_selected_data(song_id):  # Use your SQLite helper
        flash('Song not found!', 'danger')
        return redirect(url_for('mezmur'))
    
    # Validate playlist belongs to user
    playlist = Playlist.query.get(playlist_id)
    if not playlist or playlist.user_id != current_user.id:
        flash('Invalid playlist!', 'danger')
        return redirect(url_for('mezmur'))
    
    # Check if song is already in the playlist
    existing = PlaylistSong.query.filter_by(
        playlist_id=playlist_id, 
        song_id=song_id
    ).first()
    if existing:
        flash('Song already in playlist!', 'warning')
        return redirect(url_for('mezmur'))
    
    # Add to playlist
    new_entry = PlaylistSong(playlist_id=playlist_id, song_id=song_id)
    account_db.session.add(new_entry)
    account_db.session.commit()
    flash('Song added to playlist!', 'success')
    return redirect(url_for('view_playlist', playlist_id=playlist_id))

@app.route('/playlist/<int:playlist_id>/add/<int:song_id>', methods=['POST'])
@login_required
def add_to_playlist(playlist_id, song_id):
    try:
        print(f"Attempting to add song {song_id} to playlist {playlist_id}")
        
        # Retrieve song data (once)
        song_data = db.get_selected_data(song_id)
        if not song_data:
            return jsonify({
                'status': 'error',
                'message': 'Song not found'
            }), 404
        print("Song Data:", song_data)
        
        # Retrieve and validate playlist using get_or_404
        playlist = Playlist.query.get_or_404(playlist_id)
        print("Playlist:", playlist)
        
        # Ensure the playlist belongs to the current user
        if playlist.user_id != current_user.id:
            return jsonify({
                'status': 'error',
                'message': 'Invalid playlist'
            }), 403
        
        # Check if the song is already in the playlist
        if PlaylistSong.query.filter_by(playlist_id=playlist_id, song_id=song_id).first():
            return jsonify({
                'status': 'warning',
                'message': 'Song already in playlist'
            }), 409

        # Add the song to the playlist
        new_entry = PlaylistSong(playlist_id=playlist_id, song_id=song_id)
        account_db.session.add(new_entry)
        account_db.session.commit()
        
        # Create a single consolidated JSON response
        response_data = {
            'status': 'success',
            'message': 'Song added successfully',
            'playlist': {
                'id': playlist.id,
                'name': playlist.name,
                'song_count': playlist.songs.count()
            },
            'song': {
                'id': song_id,
                'title': song_data[1]  # assuming the title is in the second position
            }
        }
        
        print("Returning JSON:", response_data)
        return jsonify(response_data), 201

    except Exception as e:
        account_db.session.rollback()
        print("Error occurred:", str(e))
        return jsonify({'status': 'error', 'message': str(e)}), 500

    
# Remove a song from a playlist
@app.route('/remove_from_playlist/<int:playlist_id>/<int:song_id>', methods=['GET','POST'])
@login_required
def remove_from_playlist_form(playlist_id, song_id):
    playlist = Playlist.query.get_or_404(playlist_id)
    if playlist.user_id != current_user.id:
        flash('Access denied!', 'danger')
        return redirect(url_for('view_playlists'))
    
    entry = PlaylistSong.query.filter_by(
        playlist_id=playlist_id, 
        song_id=song_id
    ).first()
    
    if entry:
        account_db.session.delete(entry)
        account_db.session.commit()
        flash('Song removed from playlist!', 'success')
    
    return redirect(url_for('view_playlist', playlist_id=playlist_id))

@app.route('/api/playlists/<int:playlist_id>/songs/<int:song_id>', methods=['DELETE'])
@login_required
def remove_from_playlist(playlist_id, song_id):
    try:
        # Verify playlist ownership
        playlist = Playlist.query.filter_by(
            id=playlist_id,
            user_id=current_user.id
        ).first_or_404()

        # Find and delete the song relationship
        song_rel = PlaylistSong.query.filter_by(
            playlist_id=playlist_id,
            song_id=song_id
        ).first()

        if not song_rel:
            return jsonify({'status': 'error', 'message': 'Song not in playlist'}), 404

        account_db.session.delete(song_rel)
        account_db.session.commit()

        return jsonify({
            'status': 'success',
            'message': 'Song removed from playlist'
        })

    except Exception as e:
        account_db.session.rollback()
        return jsonify({'status': 'error', 'message': str(e)}), 500
    
################################PlayList end################
@app.route('/test',  methods=['GET'])
def test():
    users = [
        {
            'username':'john_doe',
            'email':'john@example.com',
            'roles':[{'name':'admin'},{'name':'editor'}]        
        },
        {
            'username':'alex_amor',
            'email':'alex@example.com',
            'roles':[{'name':'user'}]        
        }
    ]
    return render_template("user_manager.html", users=users)

@app.route('/dashboard')
@login_required
def dashboard():
    return f'hello,{current_user.username}! welcome to your dashboard' 
  


###########################Share Playlist #########################
@app.route('/playlist/shared/<int:playlist_id>')
def view_shared_playlist(playlist_id):
    print(f"DEBUG: Accessing shared playlist {playlist_id}")
    print(f"DEBUG: Current user: {current_user.username if current_user.is_authenticated else 'Anonymous'}")
    print(f"DEBUG: User roles: {[role.name for role in current_user.roles] if current_user.is_authenticated else 'None'}")
    
    try:
        playlist = Playlist.query.get(playlist_id)
        if not playlist:
            print(f"DEBUG: Playlist {playlist_id} not found in database")
            flash('Playlist not found.', 'error')
            return redirect(url_for('index'))
        
        print(f"DEBUG: Playlist found: {playlist.name}, shared: {playlist.shared}, owner: {playlist.user.username}")
        
        # Allow access if:
        # 1. User is authenticated (any logged-in user can view shared playlists)
        # 2. Playlist is marked as shared
        # 3. User is admin (admin can access any playlist)
        # 4. User is the owner (owner can always access their playlist)
        is_admin = current_user.is_authenticated and 'admin' in [role.name for role in current_user.roles]
        is_owner = current_user.is_authenticated and current_user.id == playlist.user_id
        is_authenticated = current_user.is_authenticated
        
        if not is_authenticated:
            print(f"DEBUG: Access denied - user not authenticated")
            flash('Please log in to view shared playlists.', 'warning')
            return redirect(url_for('login'))
        
        if not playlist.shared and not is_admin and not is_owner:
            print(f"DEBUG: Access denied - playlist not shared and user is not admin/owner")
            flash('This playlist is not publicly shared and you do not have permission to access it.', 'warning')
            return redirect(url_for('index'))
        
        print(f"DEBUG: Access granted - authenticated: {is_authenticated}, shared: {playlist.shared}, admin: {is_admin}, owner: {is_owner}")
        
        # Fetch song details from SQLite (`db`)
        songs = []
        for song_rel in playlist.songs.all():
            song_data = db.get_selected_data(song_rel.song_id)
            if song_data:
                songs.append(song_data)
        
        print(f"DEBUG: Rendering shared_playlist.html with {len(songs)} songs")
        return render_template('shared_playlist.html', playlist=playlist, songs=songs)
        
    except Exception as e:
        print(f"DEBUG: Exception in view_shared_playlist: {e}")
        flash('An error occurred while loading the playlist.', 'error')
        return redirect(url_for('index'))

@app.route('/playlists/discover')
def discover_playlists():
    """Display all publicly shared playlists from all users"""
    shared_playlists = Playlist.query.filter_by(shared=True).order_by(Playlist.created_at.desc()).all()
    return render_template('discover_playlists.html', playlists=shared_playlists)

@app.route('/api/playlists/clone', methods=['POST'])
@login_required
def clone_playlist():
    """Clone a public playlist to the current user's collection"""
    try:
        data = rq.get_json()
        source_playlist_id = data.get('source_playlist_id')
        new_name = data.get('name', '').strip()
        new_description = data.get('description', '').strip()
        
        if not source_playlist_id or not new_name:
            return jsonify({'status': 'error', 'message': 'Missing required fields'}), 400
        
        # Get the source playlist
        source_playlist = Playlist.query.get_or_404(source_playlist_id)
        
        # Verify it's shared
        if not source_playlist.shared:
            return jsonify({'status': 'error', 'message': 'Playlist is not publicly shared'}), 403
        
        # Check if user already has a playlist with this name
        existing = Playlist.query.filter_by(user_id=current_user.id, name=new_name).first()
        if existing:
            return jsonify({'status': 'error', 'message': 'You already have a playlist with this name'}), 409
        
        # Create new playlist
        new_playlist = Playlist(
            user_id=current_user.id,
            name=new_name,
            description=new_description,
            shared=False  # New playlist is private by default
        )
        account_db.session.add(new_playlist)
        account_db.session.flush()  # Get the ID
        
        # Copy all songs from source playlist
        source_songs = source_playlist.songs.all()
        for song_rel in source_songs:
            new_song_rel = PlaylistSong(
                playlist_id=new_playlist.id,
                song_id=song_rel.song_id
            )
            account_db.session.add(new_song_rel)
        
        account_db.session.commit()
        
        return jsonify({
            'status': 'success',
            'message': f'Playlist "{new_name}" cloned successfully',
            'playlist_id': new_playlist.id
        })
        
    except Exception as e:
        account_db.session.rollback()
        app.logger.error(f"Error cloning playlist: {str(e)}")
        return jsonify({'status': 'error', 'message': str(e)}), 500
